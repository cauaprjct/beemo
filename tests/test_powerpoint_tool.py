"""Unit tests for PowerPointTool.

Tests the PowerPoint file manipulation functionality including reading,
creating, updating, and extracting text from presentations.
"""

import os
import tempfile
from pathlib import Path

import pytest
from pptx import Presentation

from src.exceptions import CorruptedFileError
from src.powerpoint_tool import PowerPointTool


@pytest.fixture
def powerpoint_tool():
    """Fixture that provides a PowerPointTool instance."""
    return PowerPointTool()


@pytest.fixture
def temp_pptx_file():
    """Fixture that provides a temporary PowerPoint file path."""
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp_path = tmp.name
    yield tmp_path
    # Cleanup
    if os.path.exists(tmp_path):
        os.remove(tmp_path)


@pytest.fixture
def sample_pptx_file(temp_pptx_file):
    """Fixture that creates a sample PowerPoint file for testing."""
    prs = Presentation()
    
    # Add first slide with title and content
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide1 = prs.slides.add_slide(slide_layout)
    slide1.shapes.title.text = "First Slide"
    
    # Add content to first slide
    for shape in slide1.placeholders:
        if shape.placeholder_format.idx == 1:
            text_frame = shape.text_frame
            text_frame.text = "Content line 1"
            p = text_frame.add_paragraph()
            p.text = "Content line 2"
            break
    
    # Add second slide
    slide2 = prs.slides.add_slide(slide_layout)
    slide2.shapes.title.text = "Second Slide"
    
    prs.save(temp_pptx_file)
    return temp_pptx_file


class TestReadPowerPoint:
    """Tests for read_powerpoint method."""
    
    def test_read_existing_file(self, powerpoint_tool, sample_pptx_file):
        """Test reading an existing PowerPoint file."""
        result = powerpoint_tool.read_powerpoint(sample_pptx_file)
        
        assert isinstance(result, list)
        assert len(result) == 2
        
        # Check first slide
        assert result[0]["index"] == 0
        assert result[0]["title"] == "First Slide"
        assert len(result[0]["content"]) > 0
        
        # Check second slide
        assert result[1]["index"] == 1
        assert result[1]["title"] == "Second Slide"
    
    def test_read_nonexistent_file(self, powerpoint_tool):
        """Test reading a file that doesn't exist."""
        with pytest.raises(FileNotFoundError, match="PowerPoint file not found"):
            powerpoint_tool.read_powerpoint("nonexistent.pptx")
    
    def test_read_corrupted_file(self, powerpoint_tool, temp_pptx_file):
        """Test reading a corrupted PowerPoint file."""
        # Create a corrupted file
        with open(temp_pptx_file, "w") as f:
            f.write("This is not a valid PowerPoint file")
        
        with pytest.raises(CorruptedFileError, match="corrupted or invalid"):
            powerpoint_tool.read_powerpoint(temp_pptx_file)


class TestCreatePowerPoint:
    """Tests for create_powerpoint method."""
    
    def test_create_new_file(self, powerpoint_tool, temp_pptx_file):
        """Test creating a new PowerPoint file."""
        slides_data = [
            {"title": "Title 1", "content": ["Point 1", "Point 2"]},
            {"title": "Title 2", "content": ["Point A", "Point B"]}
        ]
        
        powerpoint_tool.create_powerpoint(temp_pptx_file, slides_data)
        
        # Verify file was created
        assert os.path.exists(temp_pptx_file)
        
        # Verify content
        result = powerpoint_tool.read_powerpoint(temp_pptx_file)
        assert len(result) == 2
        assert result[0]["title"] == "Title 1"
        assert result[1]["title"] == "Title 2"
    
    def test_create_empty_presentation(self, powerpoint_tool, temp_pptx_file):
        """Test creating an empty presentation."""
        powerpoint_tool.create_powerpoint(temp_pptx_file, [])
        
        assert os.path.exists(temp_pptx_file)
        result = powerpoint_tool.read_powerpoint(temp_pptx_file)
        assert len(result) == 0


class TestAddSlide:
    """Tests for add_slide method."""
    
    def test_add_slide_to_existing_file(self, powerpoint_tool, sample_pptx_file):
        """Test adding a slide to an existing presentation."""
        new_slide_data = {
            "title": "New Slide",
            "content": ["New content 1", "New content 2"],
            "notes": "These are notes"
        }
        
        powerpoint_tool.add_slide(sample_pptx_file, new_slide_data)
        
        # Verify slide was added
        result = powerpoint_tool.read_powerpoint(sample_pptx_file)
        assert len(result) == 3
        assert result[2]["title"] == "New Slide"
        assert result[2]["notes"] == "These are notes"
    
    def test_add_slide_to_nonexistent_file(self, powerpoint_tool):
        """Test adding a slide to a file that doesn't exist."""
        with pytest.raises(FileNotFoundError, match="PowerPoint file not found"):
            powerpoint_tool.add_slide("nonexistent.pptx", {"title": "Test"})


class TestUpdateSlide:
    """Tests for update_slide method."""
    
    def test_update_existing_slide(self, powerpoint_tool, sample_pptx_file):
        """Test updating an existing slide."""
        new_content = {
            "title": "Updated Title",
            "content": ["Updated content"],
            "notes": "Updated notes"
        }
        
        powerpoint_tool.update_slide(sample_pptx_file, 0, new_content)
        
        # Verify slide was updated
        result = powerpoint_tool.read_powerpoint(sample_pptx_file)
        assert result[0]["title"] == "Updated Title"
        assert result[0]["notes"] == "Updated notes"
    
    def test_update_slide_invalid_index(self, powerpoint_tool, sample_pptx_file):
        """Test updating a slide with invalid index."""
        with pytest.raises(IndexError, match="out of range"):
            powerpoint_tool.update_slide(sample_pptx_file, 999, {"title": "Test"})
    
    def test_update_slide_nonexistent_file(self, powerpoint_tool):
        """Test updating a slide in a file that doesn't exist."""
        with pytest.raises(FileNotFoundError, match="PowerPoint file not found"):
            powerpoint_tool.update_slide("nonexistent.pptx", 0, {"title": "Test"})


class TestExtractText:
    """Tests for extract_text method."""
    
    def test_extract_text_from_slides(self, powerpoint_tool, sample_pptx_file):
        """Test extracting text from all slides."""
        result = powerpoint_tool.extract_text(sample_pptx_file)
        
        assert isinstance(result, list)
        assert len(result) == 2
        
        # First slide should contain title and content
        assert "First Slide" in result[0]
        
        # Second slide should contain title
        assert "Second Slide" in result[1]
    
    def test_extract_text_nonexistent_file(self, powerpoint_tool):
        """Test extracting text from a file that doesn't exist."""
        with pytest.raises(FileNotFoundError, match="PowerPoint file not found"):
            powerpoint_tool.extract_text("nonexistent.pptx")
    
    def test_extract_text_corrupted_file(self, powerpoint_tool, temp_pptx_file):
        """Test extracting text from a corrupted file."""
        # Create a corrupted file
        with open(temp_pptx_file, "w") as f:
            f.write("This is not a valid PowerPoint file")
        
        with pytest.raises(CorruptedFileError, match="corrupted or invalid"):
            powerpoint_tool.extract_text(temp_pptx_file)
