"""PowerPoint file manipulation tool for Gemini Office Agent.

This module provides the PowerPointTool class for reading, creating, and modifying
PowerPoint files (.pptx) using the python-pptx library.
"""

import copy
import logging
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.exc import PackageNotFoundError
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from lxml import etree

from src.exceptions import CorruptedFileError
from src.logging_config import get_logger, log_file_access_error

logger = get_logger(__name__)


class PowerPointTool:
    """Tool for manipulating PowerPoint files (.pptx).
    
    Provides methods for reading, creating, updating, and adding slides
    to PowerPoint presentations using python-pptx.
    """

    def _populate_slide(self, slide, slide_data: Dict[str, Any]) -> None:
        """Popula um slide com título, conteúdo e notas.
        
        Args:
            slide: Objeto slide do python-pptx
            slide_data: Dicionário com 'title', 'content' e/ou 'notes'
        """
        # Set title if provided
        if "title" in slide_data and slide_data["title"]:
            if slide.shapes.title:
                slide.shapes.title.text = slide_data["title"]
        
        # Add content if provided
        if "content" in slide_data and slide_data["content"]:
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:  # Content placeholder
                    text_frame = shape.text_frame
                    text_frame.clear()
                    for content_item in slide_data["content"]:
                        p = text_frame.add_paragraph()
                        p.text = content_item
                    break
        
        # Add notes if provided
        if "notes" in slide_data and slide_data["notes"]:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_data["notes"]

    def read_powerpoint(self, file_path: str) -> List[Dict[str, Any]]:
        """Lê apresentação e retorna informações dos slides.
        
        Args:
            file_path: Path to the PowerPoint file to read
            
        Returns:
            List of dictionaries, each containing:
                - index: Slide index (0-based)
                - title: Slide title (if present)
                - content: List of text content from shapes
                - notes: Presenter notes (if present)
                
        Raises:
            FileNotFoundError: If the file does not exist
            CorruptedFileError: If the file is corrupted or cannot be read
        """
        logger.info(f"Reading PowerPoint file: {file_path}")
        
        path = Path(file_path)
        if not path.exists():
            log_file_access_error(logger, file_path, "arquivo não encontrado")
            raise FileNotFoundError(f"PowerPoint file not found: {file_path}")
        
        try:
            prs = Presentation(file_path)
            
            # Extract data from all slides
            slides_data = []
            for idx, slide in enumerate(prs.slides):
                slide_info = {
                    "index": idx,
                    "title": "",
                    "content": [],
                    "notes": ""
                }
                
                # Extract title (usually the first shape with text)
                if slide.shapes.title:
                    slide_info["title"] = slide.shapes.title.text
                
                # Extract content from all text shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        # Skip title if already captured
                        if shape == slide.shapes.title:
                            continue
                        slide_info["content"].append(shape.text)
                
                # Extract notes
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if notes_slide.notes_text_frame:
                        slide_info["notes"] = notes_slide.notes_text_frame.text
                
                slides_data.append(slide_info)
            
            logger.info(f"Successfully read PowerPoint file with {len(slides_data)} slides")
            return slides_data
            
        except PackageNotFoundError as e:
            logger.error(f"Corrupted PowerPoint file: {file_path} - {str(e)}")
            raise CorruptedFileError(f"PowerPoint file is corrupted or invalid: {file_path}") from e
        except Exception as e:
            logger.error(f"Error reading PowerPoint file: {file_path} - {str(e)}")
            raise CorruptedFileError(f"PowerPoint file is corrupted or invalid: {file_path} - {str(e)}") from e

    # Mapa de layouts comuns: nome amigável → índice padrão no template
    LAYOUT_MAP = {
        'title': 0,           # Title Slide (título grande + subtítulo)
        'content': 1,         # Title and Content (padrão)
        'section': 2,         # Section Header
        'two_content': 3,     # Two Content
        'comparison': 4,      # Comparison
        'title_only': 5,      # Title Only
        'blank': 6,           # Blank
    }

    def create_powerpoint(self, file_path: str, slides: List[Dict[str, Any]]) -> None:
        """Cria nova apresentação com slides fornecidos.
        
        Args:
            file_path: Path where the PowerPoint file should be created
            slides: List of dictionaries, each containing:
                - title: Slide title (optional)
                - content: List of text content items (optional)
                - notes: Presenter notes (optional)
                - layout: Layout name: 'title', 'content', 'section', 'blank', etc. (optional, default 'content')
            
        Raises:
            IOError: If the file cannot be created
        """
        logger.info(f"Creating PowerPoint file: {file_path}")
        
        Path(file_path).parent.mkdir(parents=True, exist_ok=True)
        
        try:
            prs = Presentation()
            
            # Create slides with data
            for slide_data in slides:
                layout_name = slide_data.get('layout', 'content')
                layout_idx = self.LAYOUT_MAP.get(layout_name, 1)
                try:
                    slide_layout = prs.slide_layouts[layout_idx]
                except IndexError:
                    slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                self._populate_slide(slide, slide_data)
            
            # Save the presentation
            prs.save(file_path)
            
            logger.info(f"Successfully created PowerPoint file with {len(slides)} slides")
            
        except Exception as e:
            logger.error(f"Error creating PowerPoint file: {file_path} - {str(e)}")
            raise IOError(f"Failed to create PowerPoint file: {file_path} - {str(e)}") from e

    def add_slide(self, file_path: str, slide_data: Dict[str, Any]) -> None:
        """Adiciona novo slide à apresentação.
        
        Args:
            file_path: Path to the PowerPoint file
            slide_data: Dictionary containing:
                - title: Slide title (optional)
                - content: List of text content items (optional)
                - notes: Presenter notes (optional)
            
        Raises:
            FileNotFoundError: If the file does not exist
            CorruptedFileError: If the file is corrupted
        """
        logger.info(f"Adding slide to {file_path}")
        
        path = Path(file_path)
        if not path.exists():
            log_file_access_error(logger, file_path, "arquivo não encontrado")
            raise FileNotFoundError(f"PowerPoint file not found: {file_path}")
        
        try:
            prs = Presentation(file_path)
            
            # Use title and content layout (layout index 1)
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            self._populate_slide(slide, slide_data)
            
            # Save the presentation
            prs.save(file_path)
            
            logger.info(f"Successfully added slide to presentation")
            
        except PackageNotFoundError as e:
            logger.error(f"Corrupted PowerPoint file: {file_path} - {str(e)}")
            raise CorruptedFileError(f"PowerPoint file is corrupted or invalid: {file_path}") from e
        except Exception as e:
            logger.error(f"Error adding slide to PowerPoint file: {file_path} - {str(e)}")
            raise IOError(f"Failed to add slide to PowerPoint file: {file_path} - {str(e)}") from e

    def update_slide(self, file_path: str, slide_index: int, new_content: Dict[str, Any]) -> None:
        """Atualiza conteúdo de slide específico.
        
        Args:
            file_path: Path to the PowerPoint file
            slide_index: Index of the slide to update (0-based)
            new_content: Dictionary containing:
                - title: New slide title (optional)
                - content: New list of text content items (optional)
                - notes: New presenter notes (optional)
            
        Raises:
            FileNotFoundError: If the file does not exist
            CorruptedFileError: If the file is corrupted
            IndexError: If the slide index is out of range
        """
        logger.info(f"Updating slide {slide_index} in {file_path}")
        
        path = Path(file_path)
        if not path.exists():
            log_file_access_error(logger, file_path, "arquivo não encontrado")
            raise FileNotFoundError(f"PowerPoint file not found: {file_path}")
        
        try:
            prs = Presentation(file_path)
            
            # Check if index is valid
            if slide_index < 0 or slide_index >= len(prs.slides):
                logger.error(f"Slide index {slide_index} out of range (presentation has {len(prs.slides)} slides)")
                raise IndexError(f"Slide index {slide_index} out of range (presentation has {len(prs.slides)} slides)")
            
            slide = prs.slides[slide_index]
            
            self._populate_slide(slide, new_content)
            
            # Save the presentation
            prs.save(file_path)
            
            logger.info(f"Successfully updated slide {slide_index}")
            
        except PackageNotFoundError as e:
            logger.error(f"Corrupted PowerPoint file: {file_path} - {str(e)}")
            raise CorruptedFileError(f"PowerPoint file is corrupted or invalid: {file_path}") from e
        except IndexError:
            raise
        except Exception as e:
            logger.error(f"Error updating slide in PowerPoint file: {file_path} - {str(e)}")
            raise IOError(f"Failed to update slide in PowerPoint file: {file_path} - {str(e)}") from e

    def extract_text(self, file_path: str) -> List[str]:
        """Extrai texto de todos os slides.
        
        Args:
            file_path: Path to the PowerPoint file
            
        Returns:
            List of strings, each containing all text from a slide
                
        Raises:
            FileNotFoundError: If the file does not exist
            CorruptedFileError: If the file is corrupted or cannot be read
        """
        logger.info(f"Extracting text from PowerPoint file: {file_path}")
        
        path = Path(file_path)
        if not path.exists():
            log_file_access_error(logger, file_path, "arquivo não encontrado")
            raise FileNotFoundError(f"PowerPoint file not found: {file_path}")
        
        try:
            prs = Presentation(file_path)
            
            # Extract text from all slides
            slides_text = []
            for slide in prs.slides:
                slide_text_parts = []
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text_parts.append(shape.text)
                
                # Join all text parts for this slide
                slide_text = "\n".join(slide_text_parts)
                slides_text.append(slide_text)
            
            logger.info(f"Successfully extracted text from {len(slides_text)} slides")
            return slides_text
            
        except PackageNotFoundError as e:
            logger.error(f"Corrupted PowerPoint file: {file_path} - {str(e)}")
            raise CorruptedFileError(f"PowerPoint file is corrupted or invalid: {file_path}") from e
        except Exception as e:
            logger.error(f"Error extracting text from PowerPoint file: {file_path} - {str(e)}")
            raise CorruptedFileError(f"PowerPoint file is corrupted or invalid: {file_path} - {str(e)}") from e

    def delete_slide(self, file_path: str, slide_index: int) -> None:
        """Remove um slide da apresentação.
        
        Args:
            file_path: Path to the PowerPoint file
            slide_index: Index of the slide to delete (0-based)
            
        Raises:
            FileNotFoundError: If the file does not exist
            IndexError: If the slide index is out of range
            ValueError: If trying to delete the only slide
        """
        logger.info(f"Deleting slide {slide_index} from {file_path}")
        
        path = Path(file_path)
        if not path.exists():
            log_file_access_error(logger, file_path, "arquivo não encontrado")
            raise FileNotFoundError(f"PowerPoint file not found: {file_path}")
        
        try:
            prs = Presentation(file_path)
            
            if slide_index < 0 or slide_index >= len(prs.slides):
                raise IndexError(
                    f"Slide index {slide_index} out of range "
                    f"(presentation has {len(prs.slides)} slides)"
                )
            
            if len(prs.slides) == 1:
                raise ValueError("Cannot delete the only slide in the presentation")
            
            # python-pptx não tem delete_slide nativo, precisa manipular XML
            rId = prs.slides._sldIdLst[slide_index].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[slide_index]
            
            prs.save(file_path)
            logger.info(f"Successfully deleted slide {slide_index}")
            
        except PackageNotFoundError as e:
            raise CorruptedFileError(f"PowerPoint file is corrupted or invalid: {file_path}") from e
        except (IndexError, ValueError):
            raise
        except Exception as e:
            logger.error(f"Error deleting slide: {file_path} - {str(e)}")
            raise IOError(f"Failed to delete slide: {file_path} - {str(e)}") from e

    def duplicate_slide(self, file_path: str, slide_index: int) -> None:
        """Duplica um slide existente (adicionado no final).
        
        Args:
            file_path: Path to the PowerPoint file
            slide_index: Index of the slide to duplicate (0-based)
            
        Raises:
            FileNotFoundError: If the file does not exist
            IndexError: If the slide index is out of range
        """
        logger.info(f"Duplicating slide {slide_index} from {file_path}")
        
        path = Path(file_path)
        if not path.exists():
            log_file_access_error(logger, file_path, "arquivo não encontrado")
            raise FileNotFoundError(f"PowerPoint file not found: {file_path}")
        
        try:
            prs = Presentation(file_path)
            
            if slide_index < 0 or slide_index >= len(prs.slides):
                raise IndexError(
                    f"Slide index {slide_index} out of range "
                    f"(presentation has {len(prs.slides)} slides)"
                )
            
            # Copiar XML do slide fonte
            source_slide = prs.slides[slide_index]
            slide_layout = source_slide.slide_layout
            new_slide = prs.slides.add_slide(slide_layout)
            
            # Copiar conteúdo dos shapes
            for shape in source_slide.shapes:
                el = copy.deepcopy(shape._element)
                new_slide.shapes._spTree.append(el)
            
            # Copiar notas se existirem
            if source_slide.has_notes_slide:
                src_notes = source_slide.notes_slide.notes_text_frame.text
                if src_notes:
                    new_slide.notes_slide.notes_text_frame.text = src_notes
            
            prs.save(file_path)
            logger.info(f"Successfully duplicated slide {slide_index}")
            
        except PackageNotFoundError as e:
            raise CorruptedFileError(f"PowerPoint file is corrupted or invalid: {file_path}") from e
        except IndexError:
            raise
        except Exception as e:
            logger.error(f"Error duplicating slide: {file_path} - {str(e)}")
            raise IOError(f"Failed to duplicate slide: {file_path} - {str(e)}") from e

    def add_textbox(self, file_path: str, slide_index: int, text: str,
                    left: float = 1.0, top: float = 1.0,
                    width: float = 5.0, height: float = 1.0,
                    font_size: int = 18, bold: bool = False,
                    font_color: Optional[str] = None) -> None:
        """Adiciona caixa de texto livre a um slide.
        
        Args:
            file_path: Path to the PowerPoint file
            slide_index: Index of the slide (0-based)
            text: Text content
            left: Left position in inches (default 1.0)
            top: Top position in inches (default 1.0)
            width: Width in inches (default 5.0)
            height: Height in inches (default 1.0)
            font_size: Font size in points (default 18)
            bold: Whether text is bold (default False)
            font_color: Hex color without # (optional, e.g. "FF0000")
        """
        logger.info(f"Adding textbox to slide {slide_index} in {file_path}")
        
        path = Path(file_path)
        if not path.exists():
            log_file_access_error(logger, file_path, "arquivo não encontrado")
            raise FileNotFoundError(f"PowerPoint file not found: {file_path}")
        
        try:
            prs = Presentation(file_path)
            
            if slide_index < 0 or slide_index >= len(prs.slides):
                raise IndexError(
                    f"Slide index {slide_index} out of range "
                    f"(presentation has {len(prs.slides)} slides)"
                )
            
            slide = prs.slides[slide_index]
            txBox = slide.shapes.add_textbox(
                Inches(left), Inches(top), Inches(width), Inches(height)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_size)
            p.font.bold = bold
            if font_color:
                p.font.color.rgb = RGBColor(
                    int(font_color[0:2], 16),
                    int(font_color[2:4], 16),
                    int(font_color[4:6], 16)
                )
            
            prs.save(file_path)
            logger.info(f"Successfully added textbox to slide {slide_index}")
            
        except PackageNotFoundError as e:
            raise CorruptedFileError(f"PowerPoint file is corrupted or invalid: {file_path}") from e
        except IndexError:
            raise
        except Exception as e:
            logger.error(f"Error adding textbox: {file_path} - {str(e)}")
            raise IOError(f"Failed to add textbox: {file_path} - {str(e)}") from e

    def add_table_to_slide(self, file_path: str, slide_index: int,
                           headers: List[str], rows: List[List[str]],
                           left: float = 0.5, top: float = 2.0,
                           width: float = 9.0, height: float = 1.5) -> None:
        """Adiciona tabela a um slide.
        
        Args:
            file_path: Path to the PowerPoint file
            slide_index: Index of the slide (0-based)
            headers: List of column header names
            rows: List of rows, each row is a list of cell values
            left: Left position in inches (default 0.5)
            top: Top position in inches (default 2.0)
            width: Width in inches (default 9.0)
            height: Height in inches (default 1.5)
        """
        logger.info(f"Adding table to slide {slide_index} in {file_path}")
        
        path = Path(file_path)
        if not path.exists():
            log_file_access_error(logger, file_path, "arquivo não encontrado")
            raise FileNotFoundError(f"PowerPoint file not found: {file_path}")
        
        try:
            prs = Presentation(file_path)
            
            if slide_index < 0 or slide_index >= len(prs.slides):
                raise IndexError(
                    f"Slide index {slide_index} out of range "
                    f"(presentation has {len(prs.slides)} slides)"
                )
            
            slide = prs.slides[slide_index]
            total_rows = len(rows) + 1  # +1 for header
            total_cols = len(headers)
            
            table_shape = slide.shapes.add_table(
                total_rows, total_cols,
                Inches(left), Inches(top), Inches(width), Inches(height)
            )
            table = table_shape.table
            
            # Header row
            for i, header in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = str(header)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
                        run.font.size = Pt(12)
            
            # Data rows
            for row_idx, row_data in enumerate(rows):
                for col_idx, value in enumerate(row_data):
                    if col_idx < total_cols:
                        table.cell(row_idx + 1, col_idx).text = str(value)
            
            prs.save(file_path)
            logger.info(f"Successfully added table ({total_cols}x{len(rows)}) to slide {slide_index}")
            
        except PackageNotFoundError as e:
            raise CorruptedFileError(f"PowerPoint file is corrupted or invalid: {file_path}") from e
        except IndexError:
            raise
        except Exception as e:
            logger.error(f"Error adding table to slide: {file_path} - {str(e)}")
            raise IOError(f"Failed to add table to slide: {file_path} - {str(e)}") from e

    def replace_text(self, file_path: str, old_text: str, new_text: str) -> int:
        """Busca e substitui texto em toda a apresentação.
        
        Args:
            file_path: Path to the PowerPoint file
            old_text: Text to search for
            new_text: Replacement text
            
        Returns:
            Number of replacements made
        """
        logger.info(f"Replacing '{old_text[:30]}' with '{new_text[:30]}' in {file_path}")
        
        path = Path(file_path)
        if not path.exists():
            log_file_access_error(logger, file_path, "arquivo não encontrado")
            raise FileNotFoundError(f"PowerPoint file not found: {file_path}")
        
        try:
            prs = Presentation(file_path)
            count = 0
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if old_text in run.text:
                                    run.text = run.text.replace(old_text, new_text)
                                    count += 1
                    
                    # Tabelas dentro de slides
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                for paragraph in cell.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        if old_text in run.text:
                                            run.text = run.text.replace(old_text, new_text)
                                            count += 1
            
            prs.save(file_path)
            logger.info(f"Successfully replaced {count} occurrence(s)")
            return count
            
        except PackageNotFoundError as e:
            raise CorruptedFileError(f"PowerPoint file is corrupted or invalid: {file_path}") from e
        except Exception as e:
            logger.error(f"Error replacing text: {file_path} - {str(e)}")
            raise IOError(f"Failed to replace text: {file_path} - {str(e)}") from e
